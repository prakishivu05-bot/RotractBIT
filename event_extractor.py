import os
import glob
import json
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation

# Configure Gemini API - User must set this environment variable
GOOGLE_API_KEY = os.environ.get("GEMINI_API_KEY")
if GOOGLE_API_KEY:
    genai.configure(api_key=GOOGLE_API_KEY)

# Use the Gemini Pro model
generation_config = {
  "temperature": 0.2,
  "top_p": 1,
  "top_k": 1,
  "max_output_tokens": 1024,
}

model = genai.GenerativeModel("gemini-pro", generation_config=generation_config)

def extract_text_from_pdf(filepath):
    text = ""
    try:
        reader = PdfReader(filepath)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error reading PDF {filepath}: {e}")
    return text

def extract_text_from_ppt(filepath):
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPT {filepath}: {e}")
    return text

def process_with_ai(text, filename):
    prompt = f"""
    Analyze the following text extracted from an event flyer/document ('{filename}').
    Extract the following information regarding the event and output strictly as a JSON object:
    - "id": A unique string ID based on the event name.
    - "title": Event Name
    - "date": Date of the event strictly in ISO 8601 format (YYYY-MM-DD). If no clear date, use 'YYYY-MM-DD' of the document.
    - "category": Categorize the event as strictly one of: 'Service', 'Tech', or 'Sports'. If unsure, default to 'Service'.
    - "description": A brief description of the event (1-2 sentences).
    - "sourceFile": The filename provided.

    Document Text:
    {text[:2500]} # Limiting size for API
    """
    
    if not GOOGLE_API_KEY:
        # Fallback dummy logic if NO key is provided
        return {
            "id": filename.split('.')[0].lower().replace(" ", "-"),
            "title": f"Extracted from {filename}",
            "date": "2026-04-15",
            "category": "Service",
            "description": "Auto-extracted placeholder event.",
            "sourceFile": filename
        }

    try:
        response = model.generate_content(prompt)
        response_text = response.text.strip()
        # Clean up Markdown JSON blocks if present
        if response_text.startswith("```json"):
            response_text = response_text[7:-3]
        if response_text.startswith("```"):
            response_text = response_text[3:-3]
            
        event_data = json.loads(response_text)
        return event_data
    except Exception as e:
        print(f"Error calling AI for {filename}: {e}")
        return None

def main():
    uploads_dir = "public/uploads"
    output_file = "public/events.json"
    
    os.makedirs(uploads_dir, exist_ok=True)
    
    files = glob.glob(os.path.join(uploads_dir, "*.*"))
    events = []
    
    for filepath in files:
        filename = os.path.basename(filepath)
        ext = filepath.lower().split('.')[-1]
        text = ""
        
        if ext == "pdf":
            print(f"Extracting {filename}...")
            text = extract_text_from_pdf(filepath)
        elif ext in ["ppt", "pptx"]:
            print(f"Extracting {filename}...")
            text = extract_text_from_ppt(filepath)
        else:
            continue
            
        if text.strip():
            event_data = process_with_ai(text, filename)
            if event_data:
                events.append(event_data)
                print(f"Successfully processed: {event_data.get('title')}")
                
    with open(output_file, 'w') as f:
        json.dump(events, f, indent=4)
        print(f"Saved {len(events)} events to {output_file}")

if __name__ == "__main__":
    main()
